# -*- coding: utf-8 -*-
"""RescueAI 答辩 PPT 生成脚本（沿用模板设计语言：16:9 / 深色全屏背景 / Roboto+Noto Sans SC / 白色文字+发光点缀色）"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, '.ppt_build', 'assets')
SHOTS = os.path.join(ROOT, 'screenshots')
OUT = os.path.join(ROOT, 'RescueAI_答辩PPT.pptx')

# ---------------- 设计常量（取自模板分析） ----------------
W_EMU, H_EMU = 12192000, 6858000          # 33.87 x 19.05 cm，与模板一致
W_CM, H_CM = 33.867, 19.05
F_CN = 'Noto Sans SC'
F_EN = 'Roboto'
C_WHITE = 'FEFEFE'
C_DIM = 'C9D3E0'
C_FAINT = '93A1B3'
C_ACCENT = '38BDF8'    # 青（呼应模板光带）
C_AMBER = 'F5A623'     # 琥珀（重演剧场幕间卡色）
C_RED = 'EF5B5B'
C_GREEN = '4ADE80'
C_PURPLE = 'B48BF2'
C_CARD = '0E1522'      # 卡片底
C_CARD_EDGE = '2A3A52'

# ---------------- 素材准备：裁剪长图 ----------------
def crop(name, box, out_name):
    src = os.path.join(SHOTS, name)
    dst = os.path.join(ASSETS, out_name)
    if not os.path.exists(dst):
        Image.open(src).convert('RGB').crop(box).save(dst)
    return dst

crop('t13_01_structure_real.png', (0, 0, 1440, 860), 'shot_t13_crop.png')

def asset(n):
    return os.path.join(ASSETS, n)

def shot(n):
    return os.path.join(SHOTS, n)

# ---------------- 工具函数 ----------------
prs = Presentation()
prs.slide_width = Emu(W_EMU)
prs.slide_height = Emu(H_EMU)
BLANK = prs.slide_layouts[6]

def new_slide(bg=None):
    s = prs.slides.add_slide(BLANK)
    if bg:
        s.shapes.add_picture(bg, 0, 0, width=Emu(W_EMU), height=Emu(H_EMU))
    return s

def set_run(r, text, font=F_CN, size=14, color=C_WHITE, bold=False):
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = RGBColor.from_string(color)
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', F_CN if font == F_CN else font)

def tbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """lines: list of paragraphs；每段为 list[(text, kwargs)] 或 (text, kwargs)"""
    tb = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for text, kw in para:
            set_run(p.add_run(), text, **kw)
    return tb

def rect(s, x, y, w, h, fill=C_CARD, alpha=88, edge=C_CARD_EDGE, edge_w=1.0, radius=0.06):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    sF = sh.fill._xPr.find(qn('a:solidFill'))
    srgb = sF.find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)})
    srgb.append(a)
    if edge:
        sh.line.color.rgb = RGBColor.from_string(edge)
        sh.line.width = Pt(edge_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def bar(s, x, y, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    try:
        sh.adjustments[0] = 0.5
    except Exception:
        pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(color)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def pic_fit(s, path, x, y, max_w, max_h, edge='3A4A63'):
    """等比放入框内（居中）"""
    im = Image.open(path)
    iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = iw * scale, ih * scale
    px = x + (max_w - w) / 2
    py = y + (max_h - h) / 2
    p = s.shapes.add_picture(path, Cm(px), Cm(py), Cm(w), Cm(h))
    p.line.color.rgb = RGBColor.from_string(edge)
    p.line.width = Pt(1)
    return p

def page_header(s, kicker, title, title_size=27):
    """页首：彩色小标签 + 大标题（模板风格：白字大标题居左）"""
    tbox(s, 1.5, 0.9, 28, 0.7, [
        [(kicker, dict(font=F_EN, size=12.5, color=C_ACCENT, bold=True))]])
    tbox(s, 1.5, 1.55, 30.8, 1.6, [
        [(title, dict(font=F_CN, size=title_size, color=C_WHITE, bold=True))]])
    bar(s, 1.55, 3.35, 2.6, 0.14, C_ACCENT)

def footer(s, idx=None):
    # 页码自动按当前幻灯片序号编号，保证初赛 / 复赛两种页序下 02–14 连续无重复；
    # 显式 idx 仅作为兼容保留（不应再传入硬编码值）。
    if idx is None:
        idx = len(prs.slides._sldIdLst)
    tbox(s, 1.5, 18.25, 20, 0.6, [
        [('RescueAI · Physical AI for Earthquake Rescue', dict(font=F_EN, size=10, color=C_FAINT))]])
    tbox(s, 31.6, 18.25, 1.2, 0.6, [
        [('%02d' % idx, dict(font=F_EN, size=11, color=C_FAINT, bold=True))]],
        align=PP_ALIGN.RIGHT)

def caption(s, x, y, w, text, color=C_FAINT):
    tbox(s, x, y, w, 0.55, [[(text, dict(font=F_CN, size=10, color=color))]])


def build(out_path=None, semifinal=False):
    """生成 PPT。semifinal=False 输出初赛 11 页（原行为不变）；
    semifinal=True 在核心能力与工程可靠性之后各插入两章（复赛 15 页）。"""
    OUT_ = out_path or OUT
    if semifinal:
        from _semifinal_slides import add_semifinal_slides
        _sf = dict(globals())
    # =====================================================================
    # S1 封面
    # =====================================================================
    s = new_slide(asset('bg_cover.png'))
    tbox(s, 2.0, 4.2, 29.9, 3.2, [
        [('RescueAI', dict(font=F_EN, size=76, color=C_WHITE, bold=True))]],
        align=PP_ALIGN.CENTER)
    tbox(s, 2.0, 8.3, 29.9, 1.5, [
        [('AI 地震救援指挥平台', dict(font=F_CN, size=34, color=C_WHITE, bold=True))]],
        align=PP_ALIGN.CENTER)
    tbox(s, 2.0, 10.2, 29.9, 1.1, [
        [('以 2025-03-28 缅甸 7.9 级地震为实战案例（中国测定 7.9 级 / USGS 测定 Mw 7.7）',
          dict(font=F_CN, size=16, color=C_DIM))]], align=PP_ALIGN.CENTER)
    bar(s, 15.1, 11.85, 3.7, 0.14, C_ACCENT)
    tbox(s, 2.0, 12.4, 29.9, 1.0, [
        [('双源感知 · AI 研判 · 智能救援', dict(font=F_CN, size=19, color=C_ACCENT, bold=True))]],
        align=PP_ALIGN.CENTER)
    tbox(s, 2.0, 16.7, 29.9, 0.9, [
        [('Physical AI Hackathon 2026 · XXX 团队（占位）', dict(font=F_CN, size=14, color=C_DIM))]],
        align=PP_ALIGN.CENTER)

    # =====================================================================
    # S2 为什么做 —— 震后 5 小时的信息真空
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'WHY WE BUILD', '为什么做：震后 5 小时的信息真空')
    cards = [
        (C_RED, '预警缺失', '无人收到预警', [
            '震中距中国边境约 294 km，超出既有预警覆盖',
            '缅甸无地震预警体系，全球无任何预警发出',
            '强余震 M6.7 同位置「第二击」，无从防范']),
        (C_AMBER, '官方沉默 5 小时', '通报严重滞后', [
            '官方首报 03-28 19:15 —— 震后 5 小时以上',
            '震后约 5 分钟社媒已在流传伤亡数字，信息时差巨大',
            '灾情研判只能靠碎片信息「拍脑袋」']),
        (C_PURPLE, '通讯中断', '2,500 人失联', [
            '曼德勒华文学校约 2,500 名校外学生失联',
            '设备信号中断，失联与定位无从交叉比对',
            '物资需求不明，援助只能「盲人摸象」']),
    ]
    cx = 1.5
    for color, h1, h2, items in cards:
        rect(s, cx, 4.0, 10.1, 10.9)
        bar(s, cx + 0.55, 4.65, 1.5, 0.16, color)
        tbox(s, cx + 0.55, 5.0, 9.0, 1.1,
             [[(h1, dict(font=F_CN, size=20, color=C_WHITE, bold=True))]])
        tbox(s, cx + 0.55, 6.15, 9.0, 0.9,
             [[(h2, dict(font=F_CN, size=14, color=color, bold=True))]])
        lines = [[('·  ', dict(font=F_CN, size=13, color=color, bold=True)),
                  (t, dict(font=F_CN, size=13, color=C_DIM))] for t in items]
        tbox(s, cx + 0.55, 7.35, 9.0, 6.8, lines, space_after=10, line_spacing=1.15)
        cx += 10.55
    rect(s, 1.5, 15.45, 30.8, 2.2, fill='1A1014', edge='5C3A44')
    tbox(s, 2.1, 15.75, 29.6, 1.7, [
        [('遇难 1,644 人（震后 48 小时官方通报）· 截至 4 月上旬 3,300+，数据持续更新',
          dict(font=F_CN, size=15, color=C_WHITE, bold=True))],
        [('「以下数字来自官方通报，谨以致哀。」—— 不渲染、不停顿，把时间抢回来',
          dict(font=F_CN, size=12.5, color=C_FAINT))]], space_after=5)
    footer(s)

    # =====================================================================
    # S3 产品是什么 —— 两幕式架构
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'PRODUCT OVERVIEW', '产品是什么：一场「两幕式」实战推演')
    # 第一幕卡
    rect(s, 1.5, 4.0, 15.3, 12.1)
    tbox(s, 2.05, 4.5, 14.2, 2.0, [
        [('第一幕 · 历史重演', dict(font=F_CN, size=19, color=C_AMBER, bold=True))],
        [('真实数据驱动：监测告警 → 社情爆发 → 伤亡浮现 → 救援展开，真实时间轴 ×90 压缩回放',
          dict(font=F_CN, size=12, color=C_DIM))]], space_after=4)
    pic_fit(s, shot('final_replay.png'), 2.05, 6.9, 14.2, 8.5)
    # 第二幕卡
    rect(s, 17.2, 4.0, 15.3, 12.1)
    tbox(s, 17.75, 4.5, 14.2, 2.0, [
        [('第二幕 · 数字孪生沙盘推演', dict(font=F_CN, size=19, color=C_ACCENT, bold=True))],
        [('基于真实灾情参数：空中搜索 → 结构研判 → 派遣 → 五阶段救援 → AI 参谋复盘',
          dict(font=F_CN, size=12, color=C_DIM))]], space_after=4)
    pic_fit(s, shot('final_normal.png'), 17.75, 6.9, 14.2, 8.5)
    rect(s, 1.5, 16.45, 30.8, 1.55, fill='0D1A24', edge='1F4A5E')
    tbox(s, 2.1, 16.65, 29.8, 1.2, [
        [('口播框架句：', dict(font=F_CN, size=12.5, color=C_ACCENT, bold=True)),
         ('「这是基于真实数据的数字孪生推演沙盘 —— 社情 / 余震 / 伤亡时间线为真实数据，人员与设备为推演样本。」',
          dict(font=F_CN, size=12.5, color=C_WHITE))]])
    footer(s)

    # =====================================================================
    # S4 核心能力一 · 双通道感知
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'CORE 01 · DUAL-CHANNEL SENSING', '核心能力一 · 双通道感知：仪器台网 × 社情信号')
    # 左卡：仪器侧
    rect(s, 1.5, 4.0, 14.4, 9.6)
    bar(s, 2.05, 4.55, 1.5, 0.16, C_ACCENT)
    tbox(s, 2.05, 4.9, 13.3, 1.0,
         [[('仪器侧 · USGS 真实余震目录回放', dict(font=F_CN, size=17, color=C_WHITE, bold=True))]])
    tbox(s, 2.05, 6.15, 13.3, 6.9, [
        [('·  ', dict(font=F_CN, size=13.5, color=C_ACCENT, bold=True)),
         ('主震双口径一次性标注：中国测定 7.9 级 / USGS 测定 Mw 7.7', dict(font=F_CN, size=13.5, color=C_DIM))],
        [('·  ', dict(font=F_CN, size=13.5, color=C_ACCENT, bold=True)),
         ('余震序列真实回放（us7000pn9s），M6.7 强余震同位置第二击', dict(font=F_CN, size=13.5, color=C_DIM))],
        [('·  ', dict(font=F_CN, size=13.5, color=C_ACCENT, bold=True)),
         ('震级 / 深度 / 时序逐条可溯源，标注「USGS ComCat 历史目录回放」', dict(font=F_CN, size=13.5, color=C_DIM))]],
        space_after=11, line_spacing=1.15)
    # 右卡：社情侧
    rect(s, 16.4, 4.0, 16.0, 9.6)
    bar(s, 16.95, 4.55, 1.5, 0.16, C_AMBER)
    tbox(s, 16.95, 4.9, 15.0, 1.0,
         [[('社情侧 · 微博信号瀑布流', dict(font=F_CN, size=17, color=C_WHITE, bold=True))]])
    tbox(s, 16.95, 6.15, 15.0, 6.9, [
        [('·  ', dict(font=F_CN, size=13.5, color=C_AMBER, bold=True)),
         ('数据漏斗：53,340 条原始 → 去噪 → 初筛 → AI 甄别 → 52 条精选内嵌', dict(font=F_CN, size=13.5, color=C_DIM))],
        [('·  ', dict(font=F_CN, size=13.5, color=C_AMBER, bold=True)),
         ('最早涉震帖比主震还早 1 分 46 秒 —— 词条跑赢地震波', dict(font=F_CN, size=13.5, color=C_WHITE, bold=True))],
        [('·  ', dict(font=F_CN, size=13.5, color=C_AMBER, bold=True)),
         ('峰值小时 8,122 条；NER 提取地名自动上图，情感/损毁类型自动打标', dict(font=F_CN, size=13.5, color=C_DIM))]],
        space_after=11, line_spacing=1.15)
    # 底部：社情热力截图 + 结论条
    pic_fit(s, shot('t7_06_social_heat.png'), 1.5, 13.95, 12.2, 3.9)
    caption(s, 1.5, 17.95, 12.2, '▲ 社情热度时间分布（实测截图）')
    rect(s, 14.2, 14.3, 18.2, 3.0, fill='0D1A24', edge='1F4A5E')
    tbox(s, 14.8, 14.65, 17.0, 2.4, [
        [('官方传感器沉默时，社情就是最快的传感器。', dict(font=F_CN, size=17, color=C_WHITE, bold=True))],
        [('双通道互为校验：仪器定「震」，社情定「情」。', dict(font=F_CN, size=13, color=C_DIM))]],
        space_after=6)
    footer(s)

    # =====================================================================
    # S5 核心能力二 · AI 多模态研判
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'CORE 02 · AI MULTIMODAL ASSESSMENT', '核心能力二 · AI 多模态研判：让 AI 读懂废墟')
    # 左：截图
    rect(s, 1.5, 4.0, 16.6, 12.3)
    pic_fit(s, asset('shot_t13_crop.png'), 2.05, 4.55, 15.5, 10.3)
    caption(s, 2.05, 15.1, 15.5, '▲ 真实案例：Qwen3.7-Plus 对曼德勒震后照片的结构研判（实测截图，可换图重测）')
    # 右：要点
    rect(s, 18.5, 4.0, 13.9, 12.3)
    tbox(s, 19.05, 4.5, 12.8, 1.0,
         [[('Qwen3.7-Plus · 结构研判四值', dict(font=F_CN, size=16.5, color=C_ACCENT, bold=True))]])
    tbox(s, 19.05, 5.6, 12.8, 4.6, [
        [('① 建筑图纸匹配度　② 承重构件识别', dict(font=F_CN, size=13.5, color=C_DIM))],
        [('③ 疑似生存空隙　④ 结构稳定性评估', dict(font=F_CN, size=13.5, color=C_DIM))],
        [('→ 直接生成救援建议：该不该进、从哪进、带什么装备', dict(font=F_CN, size=13, color=C_WHITE, bold=True))]],
        space_after=7, line_spacing=1.15)
    tbox(s, 19.05, 9.3, 12.8, 1.0,
         [[('Qwen3.8-Max · 三大 AI 参谋', dict(font=F_CN, size=16.5, color=C_AMBER, bold=True))]])
    tbox(s, 19.05, 10.4, 12.8, 4.4, [
        [('·  ', dict(font=F_CN, size=13.5, color=C_AMBER, bold=True)),
         ('优先级决策理由：为什么先救 TA，逐条可解释', dict(font=F_CN, size=13.5, color=C_DIM))],
        [('·  ', dict(font=F_CN, size=13.5, color=C_AMBER, bold=True)),
         ('余震参谋：强余震风险评估与作业安全建议', dict(font=F_CN, size=13.5, color=C_DIM))],
        [('·  ', dict(font=F_CN, size=13.5, color=C_AMBER, bold=True)),
         ('复盘参谋：≤150 字复盘「如果地震再来一次」', dict(font=F_CN, size=13.5, color=C_DIM))]],
        space_after=8, line_spacing=1.15)
    bar(s, 19.05, 15.15, 12.6, 0.12, C_CARD_EDGE)
    tbox(s, 19.05, 15.45, 12.8, 0.8, [
        [('AI 定位克制：辅助与提速，人做最终决策', dict(font=F_CN, size=12.5, color=C_ACCENT, bold=True))]])
    footer(s)

    # =====================================================================
    # S6 核心能力三 · 智能救援流程
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'CORE 03 · RESCUE STATE MACHINE', '核心能力三 · 智能救援：五阶段状态机，35 秒跑通')
    # 五阶段链条
    stages = [('① 开进', C_ACCENT), ('② 勘察', C_ACCENT), ('③ 语音联络', C_GREEN),
              ('④ 破拆支护', C_AMBER), ('⑤ 医疗移交', C_AMBER)]
    sx = 1.5
    for i, (name, c) in enumerate(stages):
        rect(s, sx, 4.0, 5.35, 1.5, fill='0D1A24', edge=c, edge_w=1.4)
        tbox(s, sx, 4.0, 5.35, 1.5, [[(name, dict(font=F_CN, size=15.5, color=C_WHITE, bold=True))]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 4:
            tbox(s, sx + 5.35, 4.0, 0.85, 1.5, [[('→', dict(font=F_EN, size=18, color=c, bold=True))]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        sx += 6.2
    tbox(s, 1.5, 5.85, 30.8, 0.8, [
        [('实测：派遣 → 成功营救约 35 秒全流程推进 · 阶段②可一键喊话，AI 语音联络秒回',
          dict(font=F_CN, size=12.5, color=C_DIM))]])
    # 截图区
    pic_fit(s, shot('t15_01_rescued.png'), 1.5, 6.9, 16.3, 9.3)
    caption(s, 1.5, 16.3, 16.3, '▲ 首次成功营救：S-02「已救出」（实测截图）')
    pic_fit(s, shot('t15_02_quake_pause.png'), 18.3, 6.9, 14.0, 9.3)
    caption(s, 18.3, 16.3, 14.0, '▲ 强余震自动暂停破拆，震动停止后恢复（实测截图）')
    rect(s, 1.5, 16.95, 30.8, 1.1, fill='0D1A24', edge='1F4A5E')
    tbox(s, 2.1, 17.08, 29.8, 0.9, [
        [('余震联动：', dict(font=F_CN, size=12, color=C_AMBER, bold=True)),
         ('mag≥5 自动暂停高危作业约 3s 后恢复；4.2 级不打断 —— 安全与效率由同一套状态机裁决',
          dict(font=F_CN, size=12, color=C_DIM))]])
    footer(s)

    if semifinal:
        add_semifinal_slides(_sf, 'agent')
        add_semifinal_slides(_sf, 'tools')


    # =====================================================================
    # S7 真实性与数据底座
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'DATA FOUNDATION', '真实性与数据底座：每一条都答得上来')
    rows = [
        ('主震与余震目录（时序 / 震级 / 深度）', 'USGS ComCat（事件 us7000pn9s，重演模式标注「历史目录回放」）', C_ACCENT),
        ('社情瀑布流 / 词条 TOP11 / 帖量曲线', '2025 缅甸地震微博数据集：53,340 条全库，精选 52 条内嵌', C_AMBER),
        ('伤亡阶梯 / 救援锚点（18h 医疗队 · 40h 首例）', '新华社 · 央视 · 应急管理部 · WHO 通报，时间戳交叉核验', C_RED),
        ('失联案例（2,500 名校外学生）', '官方通报与公开报道交叉核验', C_PURPLE),
        ('沙盘底图', '曼德勒真实卫星影像（含高分四号震后影像）', C_GREEN),
    ]
    ry = 4.0
    for k, v, c in rows:
        rect(s, 1.5, ry, 19.6, 2.25)
        bar(s, 2.0, ry + 0.45, 0.16, 1.35, c)
        tbox(s, 2.4, ry + 0.3, 18.4, 1.8, [
            [(k, dict(font=F_CN, size=13.5, color=C_WHITE, bold=True))],
            [(v, dict(font=F_CN, size=11.5, color=C_DIM))]], space_after=3)
        ry += 2.55
    # 右：合规卡
    rect(s, 21.6, 4.0, 10.9, 12.75)
    bar(s, 22.15, 4.55, 1.5, 0.16, C_GREEN)
    tbox(s, 22.15, 4.9, 9.8, 1.0,
         [[('匿名化与合规', dict(font=F_CN, size=17, color=C_WHITE, bold=True))]])
    tbox(s, 22.15, 6.1, 9.8, 9.5, [
        [('·  ', dict(font=F_CN, size=13, color=C_GREEN, bold=True)),
         ('依据《个人信息保护法》全量匿名化后入库', dict(font=F_CN, size=13, color=C_DIM))],
        [('·  ', dict(font=F_CN, size=13, color=C_GREEN, bold=True)),
         ('数据漏斗：53,340 → 去噪 → 初筛 → 52 条精选，全程可追溯', dict(font=F_CN, size=13, color=C_DIM))],
        [('·  ', dict(font=F_CN, size=13, color=C_GREEN, bold=True)),
         ('NER / 情感 / 损毁标注与原始帖一一对应', dict(font=F_CN, size=13, color=C_DIM))],
        [('·  ', dict(font=F_CN, size=13, color=C_GREEN, bold=True)),
         ('评委追问「数据真实性」→ 直接出示来源与交叉验证记录', dict(font=F_CN, size=13, color=C_WHITE, bold=True))]],
        space_after=10, line_spacing=1.2)
    footer(s)

    # =====================================================================
    # S8 工程可靠性 —— 演示零翻车设计
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'ENGINEERING RELIABILITY', '工程可靠性：演示零翻车设计（技术追问弹药）')
    cards8 = [
        (C_ACCENT, '三级降级链路', ['实时 API → 预缓存数据 → 预录真实响应',
                                   'AI 视觉 20s / 文本 5s 超时自动兜底', '页面降级不崩，全程无感切换']),
        (C_AMBER, '?sim=1 断网保险丝', ['0 个 AI 请求 · 0 网络依赖', '离线震情 + 全流程可演示',
                                      '可与重演模式叠加（replay=1&sim=1）']),
        (C_GREEN, '防限流与恢复', ['60s 结果缓存防限流', '代理崩溃一条命令冷启动',
                                 '断网期间瀑布流 / 地图照常']),
        (C_PURPLE, '五套自动化回归', ['Playwright 验收脚本全 PASS', '覆盖重演 / 研判 / 救援 / 复盘 / 弱网',
                                   '每次改动回归，杜绝「现场才坏」']),
    ]
    cx = 1.5
    for c, h1, items in cards8:
        rect(s, cx, 4.0, 7.55, 11.2)
        bar(s, cx + 0.5, 4.6, 1.4, 0.16, c)
        tbox(s, cx + 0.5, 4.95, 6.6, 1.4,
             [[(h1, dict(font=F_CN, size=16.5, color=C_WHITE, bold=True))]])
        lines = [[('·  ', dict(font=F_CN, size=12.5, color=c, bold=True)),
                  (t, dict(font=F_CN, size=12.5, color=C_DIM))] for t in items]
        tbox(s, cx + 0.5, 6.6, 6.6, 8.0, lines, space_after=10, line_spacing=1.2)
        cx += 7.95
    rect(s, 1.5, 15.7, 30.8, 1.9, fill='0D1A24', edge='1F4A5E')
    tbox(s, 2.1, 16.0, 29.8, 1.4, [
        [('评委问「现场断网怎么办？」', dict(font=F_CN, size=14, color=C_WHITE, bold=True)),
         (' → 一键 ?sim=1，0 请求全流程照演；这正是我们把「可靠性」当产品能力来做的原因。',
          dict(font=F_CN, size=14, color=C_DIM))]])
    footer(s)

    if semifinal:
        add_semifinal_slides(_sf, 'risk')
        add_semifinal_slides(_sf, 'metrics')


    # =====================================================================
    # S9 实战印证
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'PROVEN IN REAL RESPONSE', '实战印证：真实战例 × AI 价值映射')
    # 左：真实战例
    rect(s, 1.5, 4.0, 13.4, 12.4)
    bar(s, 2.05, 4.55, 1.5, 0.16, C_RED)
    tbox(s, 2.05, 4.9, 12.3, 1.0,
         [[('真实战例（官方通报）', dict(font=F_CN, size=17, color=C_WHITE, bold=True))]])
    facts = [
        ('18h', '中方救援队震后约 18 小时首支进入灾区（云南 37 人医疗队）'),
        ('40h', '震后约 40 小时，中方在内比都医院废墟救出首例被困者'),
        ('1亿元', '中国政府向缅甸提供 1 亿元人道主义援助'),
    ]
    fy = 6.2
    for big, txt in facts:
        tbox(s, 2.05, fy, 3.2, 2.4, [[(big, dict(font=F_EN, size=26, color=C_AMBER, bold=True))]])
        tbox(s, 5.3, fy + 0.25, 9.1, 2.2, [[(txt, dict(font=F_CN, size=12.5, color=C_DIM))]], line_spacing=1.15)
        fy += 3.1
    # 右：AI 价值映射
    rect(s, 15.3, 4.0, 17.1, 12.4)
    bar(s, 15.85, 4.55, 1.5, 0.16, C_ACCENT)
    tbox(s, 15.85, 4.9, 16.0, 1.0,
         [[('AI 价值映射：把瓶颈逐条拆掉', dict(font=F_CN, size=17, color=C_WHITE, bold=True))]])
    maps = [
        ('预警盲区（无人收到预警）', '→', '分钟级社情感知，词条跑赢地震波'),
        ('官方沉默 5 小时', '→', '社情时间线实时填补信息真空'),
        ('2,500 人失联', '→', '失联帖 × 设备信号 × 三维地图交叉比对'),
        ('物资需求不明、援助滞后', '→', '受灾人口与物资缺口动态估算'),
    ]
    my = 6.2
    for a, arrow, b in maps:
        rect(s, 15.85, my, 6.9, 2.3, fill='101826', edge='33415A', radius=0.09)
        tbox(s, 16.2, my + 0.35, 6.3, 1.7, [[(a, dict(font=F_CN, size=12, color=C_DIM))]], line_spacing=1.1)
        tbox(s, 22.85, my, 1.0, 2.3, [[(arrow, dict(font=F_EN, size=17, color=C_ACCENT, bold=True))]],
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        rect(s, 23.95, my, 7.9, 2.3, fill='0D1A24', edge='1F4A5E', radius=0.09)
        tbox(s, 24.3, my + 0.35, 7.2, 1.7, [[(b, dict(font=F_CN, size=12, color=C_WHITE, bold=True))]], line_spacing=1.1)
        my += 2.6
    footer(s)

    # =====================================================================
    # S10 未来规划
    # =====================================================================
    s = new_slide(asset('bg_content.png'))
    page_header(s, 'ROADMAP', '未来规划：弱化项全部进路线图')
    cards10 = [
        (C_ACCENT, 'ICL 地震预警接入', '成都高新减灾研究所 · API 已对接',
         '「全国地震预警」小程序同源数据；预警推送涉资质审批，当前以数据展示为主'),
        (C_AMBER, '无人机实机遥测', '接口已预留：DRONE_TELEMETRY_SCHEMA',
         '经大疆 Cloud API / MAVLink 转发，三步接入替换仿真；画面逻辑零改动'),
        (C_GREEN, '生命信号探测接入', '下一步接入：热成像 / 心跳 / 设备信号',
         '从「社情 + 仪器」扩展到「生命体征」，感知维度闭环'),
        (C_PURPLE, '多机协同编队', '多无人机 + 地面机器人协同搜索',
         '同一遥测协议横向扩展，指挥中心统一调度'),
    ]
    cx = 1.5
    for c, h1, h2, txt in cards10:
        rect(s, cx, 4.0, 7.55, 11.0)
        bar(s, cx + 0.5, 4.6, 1.4, 0.16, c)
        tbox(s, cx + 0.5, 4.95, 6.6, 1.0,
             [[(h1, dict(font=F_CN, size=16, color=C_WHITE, bold=True))]])
        tbox(s, cx + 0.5, 6.15, 6.6, 1.6,
             [[(h2, dict(font=F_EN, size=11.5, color=c, bold=True))]], line_spacing=1.1)
        tbox(s, cx + 0.5, 8.1, 6.6, 6.4,
             [[(txt, dict(font=F_CN, size=12, color=C_DIM))]], line_spacing=1.25)
        cx += 7.95
    rect(s, 1.5, 15.55, 30.8, 1.95, fill='0D1A24', edge='1F4A5E')
    tbox(s, 2.1, 15.85, 29.8, 1.5, [
        [('口播稿：', dict(font=F_CN, size=13, color=C_ACCENT, bold=True)),
         ('「当前为数字孪生仿真驱动，同一遥测协议可直接替换为实机。」—— 感知 → 分析 → 行动全链路闭环。',
          dict(font=F_CN, size=13, color=C_WHITE))]])
    footer(s)

    # =====================================================================
    # S11 结尾
    # =====================================================================
    s = new_slide(asset('bg_end.png'))
    tbox(s, 2.0, 5.2, 29.9, 2.6, [
        [('AI 不只能看数据，', dict(font=F_CN, size=40, color=C_WHITE, bold=True))],
        [('它能读懂废墟下的呼救。', dict(font=F_CN, size=40, color=C_ACCENT, bold=True))]],
        align=PP_ALIGN.CENTER, space_after=6)
    bar(s, 15.1, 10.0, 3.7, 0.14, C_ACCENT)
    tbox(s, 2.0, 10.7, 29.9, 1.0, [
        [('AI 不替代救援者，它做的只是把时间抢回来。', dict(font=F_CN, size=17, color=C_DIM))]],
        align=PP_ALIGN.CENTER)
    tbox(s, 2.0, 13.3, 29.9, 2.6, [
        [('XXX 团队（占位）', dict(font=F_CN, size=22, color=C_WHITE, bold=True))],
        [('成员 / 联系方式 / GitHub · 演示地址占位', dict(font=F_CN, size=13, color=C_FAINT))],
        [('Thank You · 致谢所有一线救援者', dict(font=F_EN, size=14, color=C_DIM))]],
        align=PP_ALIGN.CENTER, space_after=8)



    prs.core_properties.title = ('RescueAI · 复赛更新版项目方案'
                                 if semifinal else 'RescueAI · AI 地震救援指挥平台 答辩 PPT')
    prs.core_properties.author = 'XXX 团队'
    prs.save(OUT_)
    print('saved:', OUT_, 'slides =', len(prs.slides._sldIdLst))


if __name__ == '__main__':
    import sys as _sys
    if len(_sys.argv) > 1:
        build(_sys.argv[1], semifinal=True)   # 复赛版（15 页）
    else:
        build()                               # 初赛版（原输出 RescueAI_答辩PPT.pptx）
