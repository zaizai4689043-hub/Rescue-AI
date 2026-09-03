# -*- coding: utf-8 -*-
"""分析模板 PPTX：尺寸/母版/版式/配色/字体/页面结构"""
import sys
from pptx import Presentation
from pptx.util import Emu

def emu_cm(v):
    return round(v / 360000, 2)

def analyze(path):
    print('=' * 70)
    print('FILE:', path)
    prs = Presentation(path)
    print('slide size: %.2f x %.2f cm (%d x %d EMU)' % (
        emu_cm(prs.slide_width), emu_cm(prs.slide_height), prs.slide_width, prs.slide_height))
    for mi, master in enumerate(prs.slide_masters):
        print('  MASTER %d: layout=%d' % (mi, len(master.slide_layouts)))
        try:
            # theme colors
            theme = master.element.getroottree()
        except Exception:
            pass
        for li, layout in enumerate(master.slide_layouts):
            ph = [(s.placeholder_format.idx, s.placeholder_format.type, s.name) for s in layout.placeholders]
            print('    layout[%d] %r placeholders=%s' % (li, layout.name, ph))
    for si, slide in enumerate(prs.slides):
        print('-' * 60)
        print('SLIDE %d layout=%r' % (si + 1, slide.slide_layout.name))
        for sh in slide.shapes:
            t = ''
            if sh.has_text_frame:
                t = sh.text_frame.text.replace('\n', ' | ')[:90]
            fill = ''
            try:
                if sh.fill.type is not None and sh.fill.type == 1:
                    fill = 'fill=%s' % sh.fill.fore_color.rgb
            except Exception:
                pass
            font = ''
            if sh.has_text_frame and sh.text_frame.paragraphs:
                runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
                if runs:
                    r0 = runs[0]
                    fname = r0.font.name
                    fsize = r0.font.size.pt if r0.font.size else None
                    col = None
                    try:
                        if r0.font.color and r0.font.color.type is not None:
                            col = r0.font.color.rgb
                    except Exception:
                        pass
                    font = 'font=%s size=%s color=%s' % (fname, fsize, col)
            print('   [%s] pos=(%.1f,%.1f) size=(%.1f x %.1f) %s %s text=%r' % (
                sh.shape_type, emu_cm(sh.left), emu_cm(sh.top),
                emu_cm(sh.width), emu_cm(sh.height), fill, font, t))

for p in sys.argv[1:]:
    analyze(p)
