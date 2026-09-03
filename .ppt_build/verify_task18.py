# -*- coding: utf-8 -*-
"""任务 #18 回读验证：复赛产物 + 初赛回归 + 占位符扫描（只读，不改产物）"""
import re, sys
from pptx import Presentation

SF = 'outputs/rescueai_pitch/RescueAI复赛方案.pptx'
OK = True

def check(cond, msg):
    global OK
    print(('  [PASS] ' if cond else '  [FAIL] ') + msg)
    if not cond:
        OK = False

def slide_texts(prs):
    out = []
    for s in prs.slides:
        txt = []
        for sh in s.shapes:
            if sh.has_text_frame:
                txt.append(sh.text_frame.text)
        out.append('\n'.join(txt))
    return out

def footer_numbers(texts):
    """页脚页码 = 每页右下角（x>30cm 且 y>17.5cm）的两位数字"""
    nums = []
    return nums

print('== 复赛产物回读：', SF)
prs = Presentation(SF)
texts = slide_texts(prs)
check(len(prs.slides._sldIdLst) == 15, 'slides == 15（实际 %d）' % len(prs.slides._sldIdLst))

# 页脚编号：取每页最后一个匹配 ^\d{2}$ 的独立文本框（页脚页码是每页唯一两位数字框）
nums = []
per_page = []
for i, s in enumerate(prs.slides):
    found = []
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if re.fullmatch(r'\d{2}', t):
            x_cm = sh.left / 360000.0
            y_cm = sh.top / 360000.0
            if x_cm > 30 and y_cm > 17:   # 右下角 = 页脚页码位
                found.append(t)
    per_page.append(found)
    nums += found
print('  每页页脚页码：', ['-' if not f else ','.join(f) for f in per_page])
check(len(nums) == 13, '含页脚页码的页 == 13（P01 封面 / P15 结尾无编号；实际 %d）' % len(nums))
check(sorted(set(nums)) == ['%02d' % i for i in range(2, 15)], '页码集合 == {02..14}')
check(len(nums) == len(set(nums)), '页码无重复')
check(nums == sorted(nums), '页码按页序递增连续')
# 逐页定位
want = {9: '09', 10: '10', 13: '13', 14: '14'}
for p, v in want.items():
    check(per_page[p - 1] == [v], 'P%02d 页码 == %s（实际 %s）' % (p, v, per_page[p - 1]))

p12 = texts[11]
check('喊话气泡' in p12, 'P12 含「喊话气泡」')
check('0.97' in p12, 'P12 含「0.97」')
check('65ms' in p12, 'P12 含「65ms」')

p09 = texts[8]
check('52 条' in p09, 'P09 含「52 条」')
check('精选 50 条' not in p09, 'P09 无「精选 50 条」')
alltext = '\n'.join(texts)
check('精选 50 条' not in alltext and '50 条精选' not in alltext, '全篇无「50 条精选/精选 50 条」残留')

# 口径红线抽查
for kw in ['7.9', 'Mw 7.7', '1,644', '1 分 46 秒', '辅助']:
    check(kw in alltext, '红线口径仍在：「%s」' % kw)

print('\n== 初赛模式回归：')
prs1 = Presentation('RescueAI_答辩PPT.pptx')
t1 = slide_texts(prs1)
check(len(prs1.slides._sldIdLst) == 11, '初赛 slides == 11（实际 %d）' % len(prs1.slides._sldIdLst))
nums1 = []
for s in prs1.slides:
    for sh in s.shapes:
        if sh.has_text_frame and re.fullmatch(r'\d{2}', sh.text_frame.text.strip()) \
           and sh.left / 360000.0 > 30 and sh.top / 360000.0 > 17:
            nums1.append(sh.text_frame.text.strip())
check(nums1 == ['%02d' % i for i in range(2, 11)], '初赛页脚 == 02..10 连续（实际 %s）' % nums1)

print('\n== 占位符扫描（两份产物）：')
pat = re.compile(r'Lorem|TBD|待填|TODO')
for name, tt in (('复赛', texts), ('初赛', t1)):
    hits = []
    for i, t in enumerate(tt):
        for m in pat.finditer(t):
            hits.append('P%02d:%s' % (i + 1, m.group(0)))
    check(not hits, '%s 模式 Lorem|TBD|待填|TODO 零残留 %s' % (name, hits if hits else ''))

print('\nRESULT:', 'ALL PASS' if OK else 'HAS FAILURE')
sys.exit(0 if OK else 1)
